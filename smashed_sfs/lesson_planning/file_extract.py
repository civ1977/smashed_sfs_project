import io

import docx
import pptx
from pypdf import PdfReader

# Plain-text extraction only, held in memory for the life of one request -
# never written to disk or the database, matching this tool's no-upload-
# persistence design (see views.py's generate_lesson_plan).
MAX_CHARS = 20000


class UnsupportedFileError(Exception):
    pass


def extract_text(uploaded_file):
    """uploaded_file is a Django UploadedFile (request.FILES['...']).
    Returns extracted text, truncated to MAX_CHARS so a large reference
    document can't blow up the prompt sent to the AI provider."""
    name = uploaded_file.name.lower()
    data = uploaded_file.read()

    if name.endswith('.docx'):
        text = _extract_docx(data)
    elif name.endswith('.pptx'):
        text = _extract_pptx(data)
    elif name.endswith('.pdf'):
        text = _extract_pdf(data)
    else:
        raise UnsupportedFileError('Only .docx, .pptx, and .pdf files are supported.')

    return text[:MAX_CHARS]


def _extract_docx(data):
    document = docx.Document(io.BytesIO(data))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text.strip())
    return '\n'.join(parts)


def _extract_pptx(data):
    presentation = pptx.Presentation(io.BytesIO(data))
    parts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
    return '\n'.join(parts)


def _extract_pdf(data):
    reader = PdfReader(io.BytesIO(data))
    parts = [page.extract_text() or '' for page in reader.pages]
    return '\n'.join(parts)
